import re
import os
import csv
import fitz
import openpyxl
from typing import Tuple, List, Dict, Any, Generator, Optional, Set
from concurrent.futures import ThreadPoolExecutor, as_completed
from tqdm import tqdm
from functools import lru_cache

from langchain_core.documents import Document
from langchain_community.vectorstores import Milvus, FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings

from pptx import Presentation
from docx import Document as DocxDocument
from doctr.io import DocumentFile
from doctr.models import ocr_predictor
from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer

import logging
from logging.handlers import RotatingFileHandler
import sys
from datetime import datetime

def setup_logging():
    log_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'logs')
    os.makedirs(log_dir, exist_ok=True)
    
    timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
    log_file = os.path.join(log_dir, f'{timestamp}.log')
    
    logger = logging.getLogger('document_processor')
    logger.setLevel(logging.DEBUG)
    
    if logger.handlers:
        logger.handlers.clear()
    
    file_handler = RotatingFileHandler(log_file, maxBytes=10485760, backupCount=5)  # 10MB per file, 5 backups
    file_handler.setLevel(logging.DEBUG)
    
    console_handler = logging.StreamHandler(sys.stdout)
    console_handler.setLevel(logging.INFO)
    
    file_formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(filename)s:%(lineno)d - %(message)s')
    console_formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
    
    file_handler.setFormatter(file_formatter)
    console_handler.setFormatter(console_formatter)
    
    logger.addHandler(file_handler)
    logger.addHandler(console_handler)
    
    return logger

logger = setup_logging()

logger.info("Document processor module initialized")
from cohere_app.Chunking_UI import db_utility

load_dotenv()
host = os.getenv("HOST", "172.16.34.233")
port = os.getenv("PORT", "19530")
MILVUS_URL = os.getenv("MILVUS_URL", f"http://{host}:{port}")

CHUNK_SIZE = 800
OVERLAP_SIZE = 200
CHARS_PER_PAGE = 3000
MAX_WORKERS = 4  
MAX_RETRIES = 3

embeddings = HuggingFaceEmbeddings(
    model_name='/home/qa-prod/Desktop/RAG_PRODUCTION/RAG_backend/cohere_app/models/sentence_transformer',
    model_kwargs={'device': "cpu"}
)
# embeddings = SentenceTransformer('/home/qa-prod/Desktop/RAG_PRODUCTION/RAG_backend/cohere_app/models/sentence_transformer')

_ocr_model = None

def get_ocr_model():
    """Lazy load OCR model when needed"""
    global _ocr_model
    if _ocr_model is None:
        _ocr_model = ocr_predictor(pretrained=True)
    return _ocr_model

OCR_LIST = []

def extract_text_pdf(pdf_path: str) -> Tuple[List[Tuple[int, str]], str]:
    """Extract text from PDF files with improved error handling"""
    text_by_page = []
    try:
        with fitz.open(pdf_path) as pdf:
            total_pages = len(pdf)
            total_text = ""
            
            for page_num in range(min(3, total_pages)):
                page = pdf[page_num]
                text = page.get_text()
                total_text += text
            
            if not total_text.strip() and total_pages > 0:
                if pdf_path not in OCR_LIST:
                    OCR_LIST.append(pdf_path)
                return [], "OCR required - ERROR"
            
            for page_num in range(total_pages):
                page = pdf[page_num]
                text = page.get_text()
                if text.strip():  
                    text_by_page.append((page_num + 1, text))
                    
        return text_by_page, "Text extraction done"
    except fitz.FileDataError:
        if pdf_path not in OCR_LIST:
            OCR_LIST.append(pdf_path)
        return [], "PDF corrupted, OCR required - ERROR"
    except Exception as e:
        logger.error(f"Error extracting text from PDF {pdf_path}: {str(e)}")
        return [], f"Can't open the file ERROR: {str(e)}"

@lru_cache(maxsize=32)
def process_ocr_document(pdf_path: str) -> Tuple[List[Tuple[int, str]], str]:
    """Extract text using OCR from PDF with caching for repeated calls"""
    try:
        text_by_page = []
        model = get_ocr_model()
        doc = DocumentFile.from_pdf(pdf_path)
        result = model(doc)
        
        for page_num, page in enumerate(result.pages, start=1):
            page_text = ""
            words = []
            
            for block in page.blocks:
                for line in block.lines:
                    for word in line.words:
                        words.append(word)
            
            sorted_words = sorted(words, key=lambda word: (word.geometry[0][1], word.geometry[0][0]))
            page_text = " ".join([word.value for word in sorted_words])
            
            if page_text.strip():  
                text_by_page.append((page_num, page_text))
        
        if not text_by_page:
            return [], "OCR produced no text - ERROR"
            
        return text_by_page, "Text extraction done"
    except Exception as e:
        logger.error(f"Error extracting OCR text from {pdf_path}: {e}")
        return [], f"OCR FILE ERROR: {str(e)}"

def process_pptx(file_path: str) -> Tuple[List[Tuple[int, str]], str]:
    """Process PowerPoint files with improved error handling"""
    try:
        prs = Presentation(file_path)
        text_by_slide = []
        
        for i, slide in enumerate(prs.slides):
            texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
                    
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_texts = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_texts.append(cell.text.strip())
                        if row_texts:
                            texts.append(" | ".join(row_texts))
            
            slide_text = "\n".join(texts)
            if slide_text.strip():  
                text_by_slide.append((i + 1, slide_text))
                
        if not text_by_slide:
            return [], "No text content found in presentation"
            
        return text_by_slide, "Text extraction done"
    except Exception as e:
        logger.exception(f"Error processing PPTX file {file_path}")
        return [], f"PPTX FILE ERROR: {str(e)}"

def process_docx(file_path: str) -> Tuple[List[Tuple[int, str]], str]:
    """Process Word documents with improved page tracking"""
    try:
        doc = DocxDocument(file_path)
        text_by_page = []
        current_page = 1
        current_page_text = []
        char_count = 0
        
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if not para_text:
                continue
            
            is_section_break = bool(para._element.xpath('.//w:sectPr'))
            is_page_break = bool(para._element.xpath('.//w:br[@w:type="page"]'))
            
            if is_section_break or is_page_break or char_count >= CHARS_PER_PAGE:
                if current_page_text:
                    page_text = '\n'.join(current_page_text)
                    if page_text.strip():
                        text_by_page.append((current_page, page_text))
                    current_page += 1
                    current_page_text = []
                    char_count = 0
            
            current_page_text.append(para_text)
            char_count += len(para_text)
        
        if current_page_text:
            page_text = '\n'.join(current_page_text)
            if page_text.strip():
                text_by_page.append((current_page, page_text))
        
        if not text_by_page:
            return [], "No text content found in document"
            
        return text_by_page, "Text extraction done"
    except Exception as e:
        logger.exception(f"Error processing DOCX file {file_path}")
        return [], f"DOCX FILE ERROR: {str(e)}"

def process_txt(file_path: str) -> Tuple[List[Tuple[int, str]], str]:
    """Process text files with improved encoding handling"""
    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                text = f.read()
            
            if not text.strip():
                return [], "Empty text file"
                
            if len(text) > CHARS_PER_PAGE:
                pages = []
                for i in range(0, len(text), CHARS_PER_PAGE):
                    page_num = i // CHARS_PER_PAGE + 1
                    page_text = text[i:i+CHARS_PER_PAGE]
                    if page_text.strip():
                        pages.append((page_num, page_text))
                return pages, "Text extraction done"
            else:
                return [(1, text)], "Text extraction done"
        except UnicodeDecodeError:
            continue
        except Exception as e:
            logger.error(f"Error processing TXT file {file_path}: {e}")
            return [], f"TEXT FILE ERROR: {str(e)}"
    
    return [], "File encoding could not be determined"

def process_xlsx(file_path: str) -> Tuple[List[Tuple[int, str]], str]:
    """Process Excel files with improved cell value handling"""
    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        text_by_sheet = []
        
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            sheet_content = []
            has_content = False
            
            for row in worksheet.iter_rows():
                row_values = []
                for cell in row:
                    cell_value = cell.value
                    if cell_value is not None:
                        if isinstance(cell_value, (int, float)):
                            cell_value = str(cell_value)
                        elif not isinstance(cell_value, str):
                            cell_value = str(cell_value)
                            
                        cell_value = cell_value.strip()
                        if cell_value:
                            row_values.append(cell_value)
                            has_content = True
                
                if row_values:
                    sheet_content.append(",".join(row_values))
            
            if has_content:
                sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(sheet_content)
                text_by_sheet.append((sheet_name, sheet_text))
        
        workbook.close()
        
        if not text_by_sheet:
            return [], "No content found in Excel file"
            
        return text_by_sheet, "Text extraction done"
    except Exception as e:
        logger.exception(f"Error processing XLSX file {file_path}")
        return [], f"EXCEL FILE ERROR: {str(e)}"

def process_csv(file_path: str) -> Tuple[List[Tuple[int, str]], str]:
    """Process CSV files with improved dialect detection"""
    try:
        with open(file_path, 'r', newline='', encoding='utf-8') as f:
            sample = f.read(4096)
        
        dialect = csv.Sniffer().sniff(sample)
        
        with open(file_path, 'r', newline='', encoding='utf-8') as f:
            reader = csv.reader(f, dialect)
            rows = list(reader)
        
        if not rows:
            return [], "Empty CSV file"
            
        text = "\n".join([",".join(row) for row in rows if any(cell.strip() for cell in row)])
        
        if not text.strip():
            return [], "No content found in CSV file"
            
        return [(1, text)], "Text extraction done"
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', newline='', encoding='latin-1') as f:
                reader = csv.reader(f)
                rows = list(reader)
                
            text = "\n".join([",".join(row) for row in rows if any(cell.strip() for cell in row)])
            
            if not text.strip():
                return [], "No content found in CSV file"
                
            return [(1, text)], "Text extraction done"
        except Exception as e:
            logger.error(f"Error processing CSV file with latin-1 encoding {file_path}: {e}")
            return [], f"CSV FILE ERROR: {str(e)}"
    except Exception as e:
        logger.error(f"Error processing CSV file {file_path}: {e}")
        return [], f"CSV FILE ERROR: {str(e)}"

def clean_text(text: str) -> str:
    """Clean extracted text with improved whitespace and garbage handling"""
    if not text:
        return ""
        
    text = re.sub(r'\s+', ' ', text)
    
    text = re.sub(r'\.{3,}', '...', text)
    
    text = re.sub(r'([!?])\1{2,}', r'\1', text)
    
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if re.match(r'^[-=_*]{3,}$', line):
            continue
            
        char_count = len(line)
        alnum_count = sum(c.isalnum() or c.isspace() for c in line)
        if char_count > 0 and (alnum_count / char_count) < 0.5:
            continue
            
        cleaned_lines.append(line)
    
    return '\n'.join(cleaned_lines).strip()

def clean_chunk(chunk: str) -> str:
    """Clean a text chunk with improved formatting removal"""
    if not chunk:
        return ""
        
    lines = chunk.splitlines()
    cleaned_lines = []
    
    for line in lines:
        stripped_line = line.strip()
        if not stripped_line:
            continue
            
        alnum_count = sum(1 for ch in stripped_line if ch.isalnum() or ch.isspace())
        if len(stripped_line) > 0 and (alnum_count / len(stripped_line)) < 0.3:
            continue
            
        cleaned_lines.append(line)
    
    cleaned_chunk = " ".join(cleaned_lines)
    
    cleaned_chunk = re.sub(r'[\-\+\|=_*]{2,}', ' ', cleaned_chunk)
    
    cleaned_chunk = re.sub(r'\s+', ' ', cleaned_chunk)
    
    return cleaned_chunk.strip()

def read_and_split_text(text_by_page: List[Tuple[int, str]], 
                        chunk_size: int = CHUNK_SIZE, 
                        overlap_size: int = OVERLAP_SIZE) -> List[Tuple[str, int, int]]:
    """
    Create chunks using sliding window with improved sentence boundary detection
    and better page tracking.
    """
    if not text_by_page:
        return []
        
    logger.info(f"Creating chunks from {len(text_by_page)} pages...")
    
    chunks = []
    current_text = ""
    current_pages = set()
    
    for page_num, text in text_by_page:
        cleaned_text = clean_chunk(text)
        if not cleaned_text:
            continue
            
        current_text += cleaned_text + " "
        current_pages.add(page_num)
        
        while len(current_text) >= chunk_size:
            end_pos = chunk_size
            
            for punct in ['.', '!', '?']:
                pos = current_text.find(punct, end_pos - 200, end_pos + 200)
                if pos > 0:
                    end_pos = pos + 1
                    break
            
            chunk = current_text[:end_pos].strip()
            
            if chunk:
                start_page = min(current_pages) if current_pages else page_num
                end_page = max(current_pages) if current_pages else page_num
                chunks.append((chunk, start_page, end_page))
            
            overlap_pos = max(0, end_pos - overlap_size)
            current_text = current_text[overlap_pos:].strip()
            
            if overlap_pos > 0:
                if current_pages:
                    current_pages = {max(current_pages)}
    
    if current_text.strip():
        start_page = min(current_pages) if current_pages else page_num
        end_page = max(current_pages) if current_pages else page_num
        chunks.append((current_text.strip(), start_page, end_page))

    logger.info(f"Created {len(chunks)} chunks")
    return chunks

def process_document(file_path: str) -> Tuple[List[Tuple[int, str]], str]:
    """Process a document based on its file extension with retry logic"""
    if not os.path.exists(file_path):
        return [], f"File not found: {file_path}"
        
    file_ext = os.path.splitext(file_path.lower())[1]
    
    for attempt in range(MAX_RETRIES):
        try:
            if file_ext in ('.pdf', '.PDF'):
                return extract_text_pdf(file_path)
            elif file_ext == '.pptx':
                return process_pptx(file_path)
            elif file_ext == '.docx':
                return process_docx(file_path)
            elif file_ext == '.txt':
                return process_txt(file_path)
            elif file_ext == '.xlsx':
                return process_xlsx(file_path)
            elif file_ext == '.csv':
                return process_csv(file_path)
            else:
                return [], f"Unsupported file format: {file_ext}"
        except Exception as e:
            if attempt < MAX_RETRIES - 1:
                logger.warning(f"Retry {attempt+1} for {file_path}: {str(e)}")
                continue
            else:
                logger.error(f"Failed after {MAX_RETRIES} attempts: {file_path}")
                return [], f"Error processing document {file_path}: {str(e)}"

def insert_into_milvus(documents: List[Document], collection_name: str) -> bool:
    """Insert documents into Milvus with retry mechanism"""
    for attempt in range(MAX_RETRIES):
        try:
            logger.info(f"Attempting to insert {len(documents)} documents into Milvus collection '{collection_name}' (attempt {attempt+1}/{MAX_RETRIES})")
            Milvus.from_documents(
                documents,
                embeddings,
                collection_name=collection_name,
                connection_args={'uri': MILVUS_URL}
            )
            logger.info(f"Successfully inserted {len(documents)} documents into Milvus collection '{collection_name}'")
            return True
        except Exception as e:
            if attempt < MAX_RETRIES - 1:
                logger.warning(f"Milvus insertion retry {attempt+1}: {str(e)}")
                continue
            else:
                logger.error(f"Failed to insert into Milvus after {MAX_RETRIES} attempts: {str(e)}")
                raise

def process_single_file(file_path: str, collection_name: str) -> Dict[str, Any]:
    """Process a single file and insert chunks into Milvus"""
    try:
        logger.info(f"Processing file: {file_path}")
        text_by_page, message = process_document(file_path)
        
        logger.debug(f"File {file_path} processing result: {message}")
        logger.debug(f"Pages extracted: {len(text_by_page)}")
        
        if not text_by_page and "ocr" in message.lower():
            if file_path not in OCR_LIST:
                OCR_LIST.append(file_path)
                logger.info(f"Added {file_path} to OCR list")
            return {
                "file": file_path,
                "status": "OCR needed",
                "message": message,
                "success": False
            }
            
        if text_by_page and 'error' not in message.lower():
            chunks = read_and_split_text(text_by_page)
            logger.info(f"Created {len(chunks)} chunks from {file_path}")
            
            if chunks:
                documents = []
                for chunk, start_page, end_page in chunks:
                    doc = Document(
                        page_content=chunk,
                        metadata={
                            'source': file_path,
                            'page': str(start_page)
                        }
                    )
                    documents.append(doc)
                
                if documents:
                    try:
                        logger.info(f"Inserting {len(documents)} documents into Milvus for {file_path}")
                        insert_into_milvus(documents, collection_name)
                        db_utility.insert_user_access(file_path, 'YES', message, collection_name)
                        logger.info(f"Successfully inserted {file_path} into Milvus")
                        return {
                            "file": file_path,
                            "status": "Success",
                            "message": message,
                            "success": True
                        }
                    except Exception as e:
                        error_message = f"Error inserting into Milvus: {str(e)}"
                        logger.error(error_message)
                        db_utility.store_error_files_with_error(collection_name, file_path, error_message)
                        return {
                            "file": file_path,
                            "status": "Error",
                            "message": error_message,
                            "success": False
                        }
            else:
                error_msg = "No valid chunks generated"
                logger.warning(f"{error_msg} for {file_path}")
                db_utility.store_error_files_with_error(collection_name, file_path, error_msg)
                return {
                    "file": file_path,
                    "status": "Error",
                    "message": error_msg,
                    "success": False
                }
        else:
            logger.warning(f"Error processing {file_path}: {message}")
            db_utility.store_error_files_with_error(collection_name, file_path, message)
            return {
                "file": file_path,
                "status": "Error",
                "message": message,
                "success": False
            }
    except Exception as e:
        error_message = f"Unexpected error processing {file_path}: {str(e)}"
        logger.error(error_message)
        db_utility.store_error_files_with_error(collection_name, file_path, error_message)
        return {
            "file": file_path,
            "status": "Error",
            "message": error_message,
            "success": False
        }

def log_heartbeat(message: str, interval_sec: int = 300) -> None:
    """Log heartbeat to ensure process is still alive and not hung"""
    import threading
    
    logger.info(f"HEARTBEAT: {message}")
    
    threading.Timer(interval_sec, log_heartbeat, args=[message, interval_sec]).start()

def log_system_info():
    """Log system information to help with debugging"""
    import platform
    
    try:
        logger.info("-------- SYSTEM INFORMATION --------")
        logger.info(f"Python version: {platform.python_version()}")
        logger.info(f"OS: {platform.platform()}")
        logger.info("------------------------------------")
    except Exception as e:
        logger.warning(f"Could not log complete system information: {e}")

def create_langchain_documents(found_files: List[str], collection_name: str) -> Generator[Dict[str, Any], None, None]:
    """
    Process files in parallel and insert into Milvus with improved monitoring
    and resource management.
    """
    log_system_info()
    
    log_heartbeat(f"Processing collection: {collection_name}", interval_sec=300)  # 5-minute heartbeat
    
    if not found_files:
        logger.warning("No files to process")
        yield {"progress_percentage": 100, "current_progress": 0, "total_files": 0, "message": "No files to process"}
        return
        
    db_utility.create_user_access(collection_name)
    db_utility.chunking_monitor()
    db_utility.create_error_files(collection_name)
    
    try:
        processed_count = 0
        total_files = len(found_files)
        
        worker_count = min(MAX_WORKERS, total_files)
        
        logger.info(f"Starting processing {total_files} files with {worker_count} workers")
        
        logger.info(f"Processing collection: {collection_name}")
        logger.info(f"Total files to process: {total_files}")
        
        with ThreadPoolExecutor(max_workers=worker_count) as executor:
            future_to_file = {
                executor.submit(process_single_file, file, collection_name): file 
                for file in found_files
            }
            
            for future in as_completed(future_to_file):
                file = future_to_file[future]
                processed_count += 1
                
                try:
                    result = future.result()
                    progress_percentage = processed_count / total_files * 100
                    
                
                    yield {
                        "progress_percentage": progress_percentage, 
                        "current_progress": processed_count, 
                        "total_files": total_files,
                        "file": file,
                        "status": result.get("status", "Unknown")
                    }
                    
                except Exception as e:
                    logger.error(f"Error processing {file}: {str(e)}")
                    db_utility.store_error_files_with_error(collection_name, file, str(e))
                    
                    yield {
                        "progress_percentage": processed_count / total_files * 100, 
                        "current_progress": processed_count, 
                        "total_files": total_files,
                        "file": file,
                        "status": "Error",
                        "message": str(e)
                    }
    
        if OCR_LIST:
            ocr_total = len(OCR_LIST)
            logger.info(f"Processing {ocr_total} OCR files")
            
            for ocr_idx, ocr_file in enumerate(OCR_LIST):
                try:
                    text_by_page, message = process_ocr_document(ocr_file)
                    
                    if text_by_page:
                        chunks = read_and_split_text(text_by_page)
                        
                        if chunks:
                            documents = []
                            for chunk, start_page, end_page in chunks:
                                doc = Document(
                                    page_content=chunk,
                                    metadata={
                                        'source': ocr_file,
                                        'page': str(start_page)
                                    }
                                )
                                documents.append(doc)
                            
                            if documents:
                                try:
                                    insert_into_milvus(documents, collection_name)
                                    db_utility.update_ocr_status(ocr_file, collection_name)
                                    status = "Success"
                                except Exception as e:
                                    error_message = f"Error inserting OCR document into Milvus: {str(e)}"
                                    logger.error(error_message)
                                    db_utility.store_error_files_with_error(collection_name, ocr_file, error_message)
                                    status = "Error"
                        else:
                            db_utility.store_error_files_with_error(collection_name, ocr_file, "OCR produced no valid chunks")
                            status = "Error"
                    else:
                        db_utility.store_error_files_with_error(collection_name, ocr_file, message)
                        status = "Error"
                        
               
                    progress = {
                        "progress_percentage": (ocr_idx + 1) / ocr_total * 100,
                        "current_progress": ocr_idx + 1,
                        "total_files": ocr_total,
                        "file": ocr_file,
                        "status": status,
                        "message": "OCR processing"
                    }
                    logger.info(f"OCR progress: {progress['current_progress']}/{progress['total_files']} - File: {ocr_file} - Status: {status}")
                    yield progress
                    
                except Exception as e:
                    error_message = f"Error processing OCR file {ocr_file}: {str(e)}"
                    logger.error(error_message)
                    db_utility.store_error_files_with_error(collection_name, ocr_file, error_message)
                    
                    # Report OCR error progress
                    yield {
                        "progress_percentage": (ocr_idx + 1) / ocr_total * 100,
                        "current_progress": ocr_idx + 1,
                        "total_files": ocr_total,
                        "file": ocr_file,
                        "status": "Error",
                        "message": error_message
                    }
    except Exception as e:
        logger.error(f"Error in document processing pipeline: {str(e)}")
        yield {
            "progress_percentage": 100,
            "current_progress": total_files if 'total_files' in locals() else 0,
            "total_files": total_files if 'total_files' in locals() else 0,
            "status": "Error",
            "message": f"Processing pipeline error: {str(e)}"
        }

def create_faiss_index(doc_path: str, faiss_folder: str):
    text_by_page, message = process_document(doc_path)
    
    if doc_path in OCR_LIST and not text_by_page:
        text_by_page, message = process_ocr_document(doc_path)
    if not text_by_page:
        raise ValueError("failed to extract text from document: {message}")

    documents = []
    for page_num, page_text in text_by_page:
        cleaned_text = clean_text(page_text)
        doc = Document(
            page_content = cleaned_text,
            metadata={'source':doc_path, 'page':str(page_num)}
        )
        documents.append(doc)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
    split_docs = text_splitter.split_documents(documents)
    faiss_index = FAISS.from_documents(split_docs, embeddings)
    desktop_path = os.path.join(os.path.expanduser("~"), "Desktop", faiss_folder)
    os.makedirs(desktop_path, exist_ok=True)
    faiss_index.save_local(desktop_path)
    print(f"FAISS index saved to: {faiss_folder}")